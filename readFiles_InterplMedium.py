#%%
import pandas as pd
import numpy as np
import os
import pytesseract
from pdfminer.high_level import extract_text
import fitz
import io
from PIL import Image
from pptx import Presentation
import zipfile
from polyglot.detect import Detector
from generateLatexFile import *


# %%
def saveFigs(zipf, zipimgpath, outputpath, figname):
    image1 = zipf.open(zipimgpath).read()
    f = open(f'{outputpath}/{figname}.png','wb')
    f.write(image1)
    return f'{outputpath}/{figname}.png'

def get_bold_list(para):
    bold_list = []
    for run in para.runs:
        bold_list.append(run.bold)
    return bold_list
def extractFiguresTextInterplMedium(docPath, filename):

    z = zipfile.ZipFile(PATH+filename)
    all_files = z.namelist()
    images = [x for x in all_files if x.startswith('ppt/media/image6')]

    prs = Presentation(docPath+filename)
    slideEn = prs.slides[-1]
    slidePt = prs.slides[-2]

    texten = ''
    textpt = ''
    for shapes in slideEn.shapes:
        if shapes.has_text_frame:
            texten += shapes.text
    for shapes in slidePt.shapes:
        if shapes.has_text_frame:
            textpt += shapes.text

    textpt = textpt.split('última semana.', 1)[1]
    texten = texten.split('\x0b\x0b', 1)[1]

    return texten.split('\n'), textpt.split('\n'), images, z


def constructLatexFileInterpMedium(docPath, filename, outputFigure):
    texsten, texstpt, images, zipf = extractFiguresTextInterplMedium(docPath, filename)

    textpt = '\section{Meio Interplanetário} \n \subsection{Responsável: Paulo Jauer} \n \n '
    texten = '\section{Interplanetary Medium} \n \subsection{Responsible: Paulo Jauer} \n \n '
    
    outfigpath = saveFigs(zipf, images[0], f"{outputFigure}", 'figureMIIndex')

    includeFigure = "\\begin{figure}[H]\n    \\centering\n    \\includegraphics[width=14cm]{./%s}\n\\end{figure}\n \\begin{itemize}\n " % '/'.join(outfigpath.split('/')[2:]) 

    texten += includeFigure
    textpt += includeFigure
    for i in range(len(texstpt)):
        if len(texstpt[i])>0:
            textpt += '\\item ' + texstpt[i] + '\n'
    textpt += "\\end{itemize} \n"

    for i in range(len(texsten)):
        if len(texsten[i])>0:
            texten += '\\item ' + texsten[i] + '\n'
    texten += "\\end{itemize} \n"

    return textpt, texten



#%%


PATH = '/home/jose/python_projects/automatedPdfReading/data/'

filename = 'briefing_06_12_2021.pptx'


#%%
#%%

textpt, texten = constructLatexFileInterpMedium(PATH, filename, outputFigure='./latexText/figures')

generateLaTexFile(textpt, EnPt=False, outputPath='./latexText', date='25/04/2022')
# %%

# %%
