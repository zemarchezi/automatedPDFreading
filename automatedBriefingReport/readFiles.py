#%%
import re
import fitz
from automatedBriefingReport.functions import *
from pdf2image import convert_from_path
from pptx import Presentation
import zipfile
from docx import Document
# from polyglot.detect import Detector
from langdetect import detect, DetectorFactory
from pdf2image import convert_from_path
from openpyxl import load_workbook
from openpyxl_image_loader import SheetImageLoader
import string
#%%
DetectorFactory.seed = 0
########################################################################################
##  SUN
########################################################################################

def extractFiguresTextSun_0(docPath, filename, outputFigure, responsible):

    
    with fitz.open(docPath + filename) as doc:
        text = ""
        for page in doc:
            text += page.get_text()

    regexPt = r"(?<=Summary)(.|\n)*(?=Resumo)"
    regexEn = r"(?<=Resumo)(.|\n)*"

    matches = re.search(regexPt, text, re.MULTILINE)
    if matches:
        ttextPt = matches.group()
    else:
        ttextPt = ''
    matches = re.search(regexEn, text, re.MULTILINE)
    if matches:
        ttextEn = matches.group()
        ttextEn.replace('END', '')
    else:
        ttextEn = ''

    ttextEn = ttextEn.strip().replace('%', '\%').replace('\n', "\\\ ")
    ttextEn = ttextEn.replace('<', '$<$').replace('>', '$>$').replace('~', '$\sim$')
    ttextPt = ttextPt.strip().replace('%', '\%').replace('\n', "\\\ ")
    ttextPt = ttextPt.replace('<', '$<$').replace('>', '$>$').replace('~', '$\sim$')
    textpt = '\section{Sol} \n \subsection{Responsável: %s}\n\n' %(responsible) +ttextEn
    texten = '\section{Sun} \n \subsection{Responsible: %s}\n\n' %(responsible) +ttextPt


    return texten, textpt

#%%
def extractFiguresTextSun_1(docPath, filename, outputFigure, responsible):
    regexPt = r"(?i)resumo"
    regexEn = r"(?i)summary"
    
    with fitz.open(docPath + filename) as doc:
        for page in range(len(doc)):
            text = doc[page].get_text()
            textsplit = text.split("\n")
            matchespt = re.search(regexPt, text, re.MULTILINE)
            matchesen = re.search(regexEn, text, re.MULTILINE)
            if matchespt:
                pagePt = page
            elif matchesen:
                pageEn = page

        ttextEn = doc[pageEn].get_text()
        ttextPt = doc[pagePt].get_text()
    
        imagesPages = {'en':list(range(pagePt+1,pageEn)), 'pt':list(range(pageEn+1,len(doc)))}

    pages = convert_from_path(docPath + filename, 500)
    pathEn = []
    pathPt = []
    for im in imagesPages.keys():
        langPage = imagesPages[im]
        for p in range(len(langPage)):
            outimagePath = f"{outputFigure}{im}_outfileSun_{p}.jpg"
            if im == 'en':
                pathEn.append(outimagePath)
            else:
                pathPt.append(outimagePath)
            pages[langPage[p]].save(outimagePath, 'JPEG')
            # pageim = doc.loadPage(langPage[p])  # number of page
            # pix = pageim.get_pixmap()
            # output = f"{im}_outfile_{p}.png"
            # pix.save(output)

    textpt = '\section{Sol} \n \subsection{Responsável: %s}\n\n'  %(responsible)  +extractItemize(ttextPt) + '\n'
    texten = '\section{Sun} \n \subsection{Responsible: %s}\n\n'  %(responsible) +extractItemize(ttextEn) +'\n'

    
    figures = """
    \\begin{figure}[H]
        \\centering
        \\includegraphics[width=14cm]{./%s}
    \\end{figure} \n \n
    """
    for s in range(len(pathEn)):
        ouIm = '/'.join(pathEn[s].split('/')[2:])
        texten += figures % (ouIm)

    for s in range(len(pathPt)):
        ouIm = '/'.join(pathPt[s].split('/')[2:])
        textpt += figures % (ouIm)

    return texten, textpt
########################################################################################



########################################################################################
##             Interplanetary Medium
########################################################################################


def extractFiguresTextInterplMedium(docPath, filename):

    z = zipfile.ZipFile(docPath+filename)
    all_files = z.namelist()
    images = [x for x in all_files if x.startswith('ppt/media/image6')]

    prs = Presentation(docPath+filename)
    slideEn = prs.slides[-1]
    slidePt = prs.slides[-2]

    texten = ''
    textpt = ''
    for shapes in slideEn.shapes:
        if shapes.has_text_frame:
            texten += shapes.text
    for shapes in slidePt.shapes:
        if shapes.has_text_frame:
            textpt += shapes.text

    splitPt = textpt.split('ltima semana.', 1)
    splitEn = texten.split('\x0b\x0b', 1)
    patternEn = r"(?i)past.*week\."
    if len(splitEn) < 2:
        splitEn = texten.split('\x0b', 1)
        if len(splitEn) < 2:
            splitEn = re.split(patternEn, aas[1])
    textpt = splitPt[1]
    texten = splitEn[1]

    return texten.split('\n'), textpt.split('\n'), images, z


def constructLatexFileInterpMedium(docPath, filename, outputFigure, responsible):
    texsten, texstpt, images, zipf = extractFiguresTextInterplMedium(docPath, filename)

    textpt = '\section{Meio Interplanetário} \n \subsection{Responsável: %s} \n \n ' %(responsible) 
    texten = '\section{Interplanetary Medium} \n \subsection{Responsible: %s} \n \n ' %(responsible) 
    
    outfigpath = saveFigs(zipf, images[0], f"{outputFigure}", 'figureMIIndex', crop=False)

    includeFigure = "\\begin{figure}[H]\n    \\centering\n    \\includegraphics[width=14cm]{./%s}\n\\end{figure}\n \\begin{itemize}\n " % '/'.join(outfigpath.split('/')[2:]) 

    texten += includeFigure
    textpt += includeFigure
    for i in range(len(texstpt)):
        if len(texstpt[i])>0:
            textpt += '\\item ' + texstpt[i] + '\n'
    textpt += "\\end{itemize} \n"

    for i in range(len(texsten)):
        if len(texsten[i])>0:
            texten += '\\item ' + texsten[i] + '\n'
    texten += "\\end{itemize} \n"

    return texten, textpt
#%%

########################################################################################
##             Radiation Belts
########################################################################################



def extractFiguresTextRadBelts(docPath, filename):

    z = zipfile.ZipFile(docPath+filename)
    all_files = z.namelist()
    images = [x for x in all_files if x.startswith('word/media/')]

    document = Document(docPath+filename)
    texs= []
    for para in document.paragraphs:
        texs.append(para.text)

    dictsposition = {}
    fig = 0
    for n, i in enumerate(texs):
        if i.startswith('Fig'):
            fig += 1
            dictsposition[f"Figure_{fig}"] = {'leg': i.split(':')[1],
                                              'fig': images[fig-1]}
        if i.startswith('Summary') or i.startswith('Resumo'):
            dictsposition[f"locSummary"] = n

    return dictsposition, texs, z


def constructLatexFileRadBelts(docPath, filename, outputFigure, responsible):
    dictsposition, texst, zipf = extractFiguresTextRadBelts(docPath, filename)
    keys = list(dictsposition.keys())
    # detector = Detector(dictsposition[keys[0]]['leg'])
    detector = detect(dictsposition[keys[0]]['leg'])
    # if detector.language.code == 'pt':
    #     text = '\section{Cinturões de Radiação} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    # if detector.language.code == 'en':
    #     text = '\section{Radiation Belts} \n \subsection{Responsible: %s} \n \n'  %(responsible) 
    if detector == 'pt':
        text = '\section{Cinturões de Radiação} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    if detector == 'en':
        text = '\section{Radiation Belts} \n \subsection{Responsible: %s} \n \n'  %(responsible) 
    keys = list(dictsposition.keys())
    
    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                             \\caption{%s}
                        \\end{figure}\n
                     """

    
    for i in range(len(keys)):
        if keys[i].startswith("Fig"):

            if keys[i].endswith("1"):
                crops = [300,1000,250,1750]
            if keys[i].endswith("2"):
                crops = [190, 1000,250,1750]
            outfigpath = saveFigs(zipf, dictsposition[keys[i]]['fig'], f"{outputFigure}", f'figureRadBelts_{i}', crops)

            text += includeFigure % ('/'.join(outfigpath.split('/')[2:]), dictsposition[keys[i]]['leg'])

    
    textSum = '\n\n'.join(texst[dictsposition['locSummary']+1:])
    text += textSum


    return text



########################################################################################
##             ULF Waves
########################################################################################


def extractFiguresTextULF(docPath, filename, outputFigure, responsible):
    ttextEn = ''
    ttextPt = ''
    with fitz.open(docPath + filename) as doc:
        for page in range(-2,0):
            text = doc[page].get_text()
            detector = detect(text)
            # if detector.language.code == 'pt':
            #     ttextPt += text
            # if detector.language.code == 'en':
            #     ttextEn += text
            if detector == 'pt':
                ttextPt += text
            if detector == 'en':
                ttextEn += text

    ttextEn += ' '.join(ttextEn.split('\n')[:-3])
    ttextPt += ' '.join(ttextPt.split('\n')[:-3])
    


    pages = convert_from_path(docPath + filename, 500)


    imagesDict = {"FigISLL": {'page': 1,
                              'leg_pt': """a) sinal do campo magnético total 
                              medido na Estação ISLL da rede CARISMA em cinza, 
                              junto com a flutuação na faixa de Pc5 em preto. b) 
                              Espectro de potência wavelet do sinal filtrado. c) 
                              Média da potência espectral nas faixas de 2 a 10 minutos 
                              (ondas ULF).""",
                              'leg_en': """a) signal of the total magnetic 
                              field measured in the ISLL Station of the CARISMA 
                              network in gray, together with the fluctuation in the 
                              range of Pc5 in black. b) Wavelet power spectrum of the 
                              filtered signal. c) Average spectral power in the ranges 
                              from 2 to 10 minutes (ULF waves).""", 
                              'crop': [260,1650,200,2900]},
                  "FigEmbrace": {'page': 2,
                              'leg_pt': """a) sinal do campo magnético total medido 
                              na Estação SMS da rede EMBRACE em cinza, junto com a 
                              flutuação na faixa de Pc5 em preto. b) Espectro de potência 
                              wavelet do sinal filtrado. c) Média da potência espectral nas 
                              faixas de 2 a 10 minutos (ondas ULF).""",
                              'leg_en': """a) signal of the total magnetic field 
                              measured in the EMBRACE network in gray, together with
                               the fluctuation in the range of Pc5 in black. b)
                                Wavelet power spectrum of the filtered signal. c) 
                                Average spectral power in the ranges from 2 to 10
                                 minutes (ULF waves).""", 
                              'crop': [260,1650,200,2900]},
                  "FigGOES": {'page': -3,
                              'leg_pt': """a) sinal do campo magnético total medido pelo 
                              satélite GOES 16, junto com a flutuação na faixa de Pc5 
                              em preto. b) Espectro de potência wavelet do sinal 
                              filtrado. c) Média da potência espectral nas faixas 
                              de 2 a 10 minutos (ondas ULF).""",
                              'leg_en': """a) signal of the total magnetic field 
                              measured by the GOES 16 satellite, together with the 
                              fluctuation in the range of Pc5 in black. b) Wavelet 
                              power spectrum of the filtered signal. c) Average 
                              spectral power in the ranges from 2 to 10 minutes 
                              (ULF waves).""", 
                              'crop': [260,1650,200,2900]},
                  }
    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                             \\caption{%s}
                        \\end{figure}\n
                     """
    keys = list(imagesDict.keys())
    textpt = '\section{Ondas ULF} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    texten = '\section{ULF Waves} \n \subsection{Responsible: %s} \n \n'  %(responsible)  

    for i in range(len(keys)):
        crops = imagesDict[keys[i]]['crop']
        outfigpath = saveFigsPdf(pages, imagesDict[keys[i]]['page'], 
                             f"{outputFigure}", f'figureULF_{i}', 
                             crops)
        
        textpt += includeFigure % ('/'.join(outfigpath.split('/')[2:]), imagesDict[keys[i]]['leg_pt'])
        texten += includeFigure % ('/'.join(outfigpath.split('/')[2:]), imagesDict[keys[i]]['leg_en'])

    textpt += ttextPt
    texten += ttextEn

    return texten, textpt


########################################################################################
##             EMIC Waves
########################################################################################


def extractFiguresTextEMIC(docPath, filename, outputFigure, responsible):
    textPt = ''
    textEn = ''
    

    pages = convert_from_path(docPath + filename, 500)

    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                        \\end{figure}\n
                     """
    textpt = '\section{Ondas EMIC} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    texten = '\section{EMIC Waves} \n \subsection{Responsible: %s} \n \n'  %(responsible) 
    for i in range(1,len(pages)):
        outfigpath = saveFigsPdf(pages, i, 
                             f"{outputFigure}", f'figureEMIC_{i}', crop=False)
        
        textpt += includeFigure % ('/'.join(outfigpath.split('/')[2:]))
        texten += includeFigure % ('/'.join(outfigpath.split('/')[2:]))

    textPt += textpt
    textEn += textpt

    return textEn, textPt


####################################################
########################################################################################
##             Geomagnetism
########################################################################################


def extractFiguresTextGeomag(docPath, filename):

    wb = load_workbook(docPath+filename)
    

    textpt = []
    texten = []

    sheeten = wb['Plan3en']
    sheetpt = wb['Plan4pt']
    for r in range(sheeten.max_row):
        cell_obj = sheeten.cell(row = r+1, column = 1)
        texten.append(cell_obj.value)

    for r in range(sheetpt.max_row):
        cell_obj = sheetpt.cell(row = r+1, column = 1)
        textpt.append(cell_obj.value)
    # Put your sheet in the loader
    sheet = wb['Plan1']
    image_loader = SheetImageLoader(sheet)

    cellFigure = []
    for nn in range(1,80):
        for ll in list(string.ascii_uppercase):
            if image_loader.image_in(f'{ll}{nn}'):
                cellFigure.append(f"{ll}{nn}")

    

    return cellFigure, textpt, texten


def constructLatexFileGeomag(docPath, filename, outputFigure, responsible):
    cellFigure, tpt, ten = extractFiguresTextGeomag(docPath, filename)

    textpt = '\section{Geomagnetismo} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    texten = '\section{Geomagnetism} \n \subsection{Responsible: %s} \n \n'  %(responsible) 


    
    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                        \\end{figure}\n
                     """
    wb = load_workbook(docPath+filename)
    sheet = wb['Plan1']
    image_loader = SheetImageLoader(sheet)
    for nn, im in enumerate(cellFigure):
        image = image_loader.get(im)
        outfigpath = saveFigsGeo(image, f"{outputFigure}", f'figureGeomag_{nn}', crop=False)
        texten += includeFigure % ('/'.join(outfigpath.split('/')[2:]))
        textpt += includeFigure % ('/'.join(outfigpath.split('/')[2:]))
    
    textpt += "\\begin{itemize} \n" 
    texten += "\\begin{itemize} \n" 
    for i in tpt:
        if not i.startswith('-'):
            textpt += "\\item " + i + "\n"
        else:
            textpt += "\\begin{itemize} \n \\item " + i + "\n \\end{itemize} \n"

    for i in ten:
        if not i.startswith('-'):
            texten += "\\item " + i + "\n"
        else:
            texten += "\\begin{itemize} \n \\item " + i + "\n \\end{itemize} \n"



    textpt += "\\end{itemize} \n"
    texten += "\\end{itemize} \n"
    # text += textSum


    return texten, textpt


####################################################
########################################################################################
##             Ionosphere
########################################################################################


def extractFiguresTextIonosphere(docPath, filename):

    z = zipfile.ZipFile(docPath+filename)
    all_files = z.namelist()
    images = [x for x in all_files if x.startswith('word/media/')]

    document = Document(docPath+filename)
    texst = list()
    bolds = []
    for para in document.paragraphs:
        bold_list = get_bold_list(para)
        if True in bold_list:
            if len(para.text) > 0:
                bolds.append(para.text)
        bolds = list(dict.fromkeys(bolds))
        if len(para.text) > 0:
            texst.append(para.text)
    bold_list = []
    for run in para.runs:
        bold_list.append(run.bold)

    dictsposition = {}
    for nnb, bol in enumerate(bolds):
        dictsposition[bol] = {}
        dictsposition[bol]['img'] = images[nnb] 
        for nn, te in enumerate(texst):
            if te.startswith(bol):
                detector = detect(texst[nn+1])
                # if detector.language.code == 'pt':
                #     dictsposition[bol]['pt'] = nn
                # if detector.language.code == 'en':
                #     dictsposition[bol]['en'] = nn
                if detector == 'pt':
                    dictsposition[bol]['pt'] = nn
                if detector == 'en':
                    dictsposition[bol]['en'] = nn

    return dictsposition, texst, z


def constructLatexFileIonosphere(docPath, filename, outputFigure, responsible):
    dictsposition, texst, zipf = extractFiguresTextIonosphere(docPath, filename)
    keys = list(dictsposition.keys())
    textpt = '\section{Ionosfera} \n \subsection{Responsável: %s} \n \n'  %(responsible) 
    texten = '\section{Ionosphere} \n \subsection{Responsible: %s} \n \n'  %(responsible) 
    for i in range(len(keys)):
        figname = keys[i].replace(' ', '').replace(':','')
        outfigpath = saveFigs(zipf, dictsposition[keys[i]]['img'], f"{outputFigure}", figname, crop=False)
        textpt += '\\textbf{%s}\n\n \\begin{itemize}\n' % (keys[i])
        texten += '\\textbf{%s}\n\n \\begin{itemize}\n' % (keys[i])
        if i+1 < len(keys):
            textitempt = ''
            for ti in texst[dictsposition[keys[i]]['pt']+1:dictsposition[keys[i+1]]['pt']]:
                textitempt += '\\item ' + ti + '\n'
            textitemen = ''
            for ti in texst[dictsposition[keys[i]]['en']+1:dictsposition[keys[i+1]]['en']]:
                textitemen += '\\item ' + ti + '\n'
        else:
            textitempt = ''
            for ti in texst[dictsposition[keys[i]]['pt']+1:dictsposition[keys[0]]['en']]:
                textitempt += '\\item ' + ti + '\n'
            textitemen = ''
            for ti in texst[dictsposition[keys[i]]['en']+1:]:
                textitemen += '\\item ' + ti + '\n'

        includeFigure = "\\begin{figure}[H]\n    \\centering\n    \\includegraphics[width=14cm]{./%s}\n\\end{figure}\n" % '/'.join(outfigpath.split('/')[2:]) 

        textpt += f'{textitempt}'+ '\\end{itemize}\n'
        textpt += f"{includeFigure}\n"
        texten += f'{textitemen}' + '\\end{itemize}\n'
        texten += f"{includeFigure}\n"

    return texten, textpt


########################################################################################
##             Ionosphere Scintilation S4
########################################################################################


def extractFiguresTextScint(docPath, filename, outputFigure, responsible):
    regexPt = r"(?i)resumo"
    regexEn = r"(?i)summary"
    
    with fitz.open(docPath + filename) as doc:
        for page in range(len(doc)):
            text = doc[page].get_text()
            textsplit = text.split("\n")
            matchespt = re.search(regexPt, text, re.MULTILINE)
            matchesen = re.search(regexEn, text, re.MULTILINE)
            if matchespt:
                pagePt = page
            elif matchesen:
                pageEn = page

        ttextEn = re.split('\n\s+',doc[pageEn].get_text())[-2]
        ttextPt = re.split('\n\s+', doc[pagePt].get_text())[-2]
    
        imagesPages = {'en':list(range(pagePt+1,pageEn)), 'pt':list(range(pageEn+1,len(doc)))}

    pages = convert_from_path(docPath + filename, 500)
    pathEn = []
    pathPt = []
    for im in imagesPages.keys():
        langPage = imagesPages[im]
        for p in range(len(langPage)):
            outimagePath = f"{outputFigure}{im}_outfileScint_{p}.jpg"
            if im == 'en':
                pathEn.append(outimagePath)
            else:
                pathPt.append(outimagePath)
            pages[langPage[p]].save(outimagePath, 'JPEG')
            # pageim = doc.loadPage(langPage[p])  # number of page
            # pix = pageim.get_pixmap()
            # output = f"{im}_outfile_{p}.png"
            # pix.save(output)

    textpt = '\section{Cintilação} \n \subsection{Responsável: %s} \n \n'  %(responsible)  +ttextPt + '\n'
    texten = '\section{Scintilation} \n \subsection{Responsible: %s} \n \n'  %(responsible)  +ttextEn +'\n'

    print(pathEn)
    
    figures = """
    \\begin{figure}[H]
        \\centering
        \\includegraphics[width=14cm]{./%s}
    \\end{figure} \n \n
    """
    for s in range(len(pathEn)):
        ouIm = '/'.join(pathEn[s].split('/')[2:])
        texten += figures % (ouIm)

    for s in range(len(pathPt)):
        ouIm = '/'.join(pathPt[s].split('/')[2:])
        textpt += figures % (ouIm)

    return texten, textpt


########################################################################################
##             All-Sky Imager
########################################################################################




def extractFiguresTextImager(docPath, filename, outputFigure, responsible):

    pages = convert_from_path(docPath + filename, 500)


    textDict = {"Tabela": {'page': 0, 
                              'crop': [1500,4000,250,3800]},
                  "imager": {'page': 1,
                              'regex_en': r"(?<=Remarks.)(.|\n)*",
                              'regex_pt': r"(?<=Observacoes.)(.|\n)*(?=Remarks)"},
                  "tec": {'page': 2,
                              'regex_en': r"(?<=Remarks.)(.|\n)*",
                              'regex_pt': r"(?<=Observacoes.)(.|\n)*(?=Remarks)"}
                  }
    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                        \\end{figure}\n
                     """
    keys = list(textDict.keys())
    # ttextEn = ''
    # ttextPt = ''
    ttextPt = '\section{Imageador All-Sky} \n \subsection{Responsável: %s} \n \n'  %(responsible)
    ttextEn = '\section{All-Sky Imager} \n \subsection{Responsible: %s} \n \n'  %(responsible)

    
    for i in range(len(keys)):
        if 'crop' in list(textDict[keys[i]].keys()):
            crops = textDict[keys[i]]['crop']
            outfigpath = saveFigsPdf(pages, textDict[keys[i]]['page'], 
                                f"{outputFigure}", f'figureImager_{i}', 
                                crops)
            ttextEn += includeFigure % ('/'.join(outfigpath.split('/')[2:]))
            ttextPt += includeFigure % ('/'.join(outfigpath.split('/')[2:]))
            # includeFigure % ('/'.join(outfigpath.split('/')[2:]))
        else:
            if keys[i] == 'tec':
                ttextPt += 'TEC \n'
                ttextEn += 'TEC \n'
            doc = fitz.open(docPath + filename)
            text = doc[textDict[keys[i]]['page']].get_text()
            text = re.sub('[\¸\˜\´\❘❙❚]+', '', text)
            matches = re.search(textDict[keys[i]]['regex_en'], text, re.MULTILINE)
            
            if matches:
                ttextEn += extractItemize_Imager(matches.group())
            matches = re.search(textDict[keys[i]]['regex_pt'], text, re.MULTILINE)
            if matches:
                ttextPt += extractItemize_Imager(matches.group())

    return ttextEn, ttextPt


########################################################################################
##             ROTI Plasma Bubble
########################################################################################

def extractFiguresTextRoti(docPath, filename):

    z = zipfile.ZipFile(docPath+filename)
    all_files = z.namelist()
    image = [x for x in all_files if x.startswith('word/media/')]
    images = []
    for im in image:
        img1 = z.open(im).read()
        img = cv2.imdecode(np.frombuffer(img1, np.uint8),3)
        if img.shape[0]>120:
            images.append(im)

    document = Document(docPath+filename)
    texs = []
    textos = []
    for para in document.paragraphs:
        texs.append(para.text)
        textos.append(para.text)
    removeList = ['Ionos', 'Carol', 'Summa', 'Resum', 'Figu', 'Tab']
    for ii in removeList:
        for nn, jj in enumerate(texs):
            if jj.startswith(ii):
                texs.pop(nn)

    indexRef = [i for i in range(len(texs)) if texs[i].startswith('Refe')]
    # print(indexRef)
    if len(indexRef) > 0:
        texs = texs[0:indexRef[0]]

    tabs = extractTables(document)
    dictsposition = {}
    fig = 0
    tab = 0
    for n, i in enumerate(textos):
        if i.startswith('Fig'):
            fig += 1
            ll = i
            if fig-1 < len(images):
                ims = images[fig-1]
            else:
                ims = images[-1]
            dictsposition[f"Figure_{fig}"] = {'leg': ll,
                                              'fig': ims}
        if i.startswith('Tab'):
            tab+=1
            ll = i
            if tab-1 < len(tabs):
                tbs = tabs[tab-1]
            else:
                tbs = tabs[-1]
            dictsposition[f"Table_{tab}"] = {'leg': i,
                                             'tab': tbs}

    return dictsposition, texs, z


def constructLatexFileRoti(docPath, filename, outputFigure, responsible):
    dictsposition, texst, zipf = extractFiguresTextRoti(docPath, filename)

    textpt = '\section{ROTI} \n \subsection{Responsável: %s} \n \n'  %(responsible)
    texten = '\section{ROTI} \n \subsection{Responsible: %s} \n \n'  %(responsible)

    for i in texst:
        if len(i)>1:
            detector = detect(i)
            # if detector.language.code == 'en':
            #     texten += i + '\n\n'
            # if detector.language.code == 'pt':
            #     textpt += i + '\n\n'
            if detector == 'en':
                texten += i + '\n\n'
            if detector == 'pt':
                textpt += i + '\n\n'
    
    # if detector.language.code == 'pt':
    #     text = 
    # if detector.language.code == 'en':
    #     text = '\section{ROTI} \n \subsection{Responsible: Carolina de Sousa do Carmo} \n \n'
    keys = list(dictsposition.keys())
    
    includeFigure = """\\begin{figure}[H]\n    
                        \\centering\n   
                             \\includegraphics[width=14cm]{./%s}\n
                             \\caption{%s}
                        \\end{figure}\n
                     """

    for i in range(len(keys)):
            
        if keys[i].startswith("Tab"):
            tbs = dictsposition[keys[i]]['tab'] % (dictsposition[keys[i]]['leg'].split('–')[-1])
            if dictsposition[keys[i]]['leg'].startswith('Table'):
                texten += tbs
            if dictsposition[keys[i]]['leg'].startswith('Tabela'):
                textpt += tbs

        if keys[i].startswith("Fig"):

            outfigpath = saveFigs(zipf, dictsposition[keys[i]]['fig'], f"{outputFigure}", f'figureROTI_{i}', crop=False)
            
            if dictsposition[keys[i]]['leg'].startswith('Figure'):
                texten += includeFigure % ('/'.join(outfigpath.split('/')[2:]), dictsposition[keys[i]]['leg'].split('–')[-1])
            if dictsposition[keys[i]]['leg'].startswith('Figura'):
                textpt += includeFigure % ('/'.join(outfigpath.split('/')[2:]), dictsposition[keys[i]]['leg'].split('–')[-1])

    # textSum = '\n\n'.join(texst[dictsposition['locSummary']+1:])
    # text += textSum


    return texten, textpt