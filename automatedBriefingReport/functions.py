import re
import io
import numpy as np
import pandas as pd
import cv2
from io import BytesIO
from openpyxl import load_workbook
from openpyxl_image_loader import SheetImageLoader
from pptx import Presentation
import fitz
from docx import Document
from openpyxl import load_workbook

def extractItemize(text):
    dictTex = {}
    for nn, i in enumerate(text.split('\n')[1:]):
        i = re.sub('[\¸\˜\´]+', '', i)

        if i.startswith("WSA") or i.startswith("EMC") or i.startswith("CME"):
            dictTex[i] = []
            ns = i
        else:
            dictTex[ns].append(i) 
    dics = {}
    for kk in dictTex.keys():
        tts = dictTex[kk]
        tts = ' '.join(tts).split('●')
        tts = [i for i in tts if len(i)>0]
        dics[kk] = tts

    texs = """\\begin{itemize} \n """

    for i in dics.keys():
        texs += '\\item ' + i +'\n'
        if len(dics[i]) > 0:
            for j in dics[i]:
                texs += """\\begin{itemize} \n """ + '\\item' + j + '\n'
        
            texs += """\\end{itemize} \n """
    
    texs += """\\end{itemize} \n """

    return texs


def saveFigs(zipf, zipimgpath, outputpath, figname, crop):
    image1 = zipf.open(zipimgpath).read()
    img = cv2.imdecode(np.frombuffer(image1, np.uint8),3)
    if crop:
        imgCropped = img[crop[0]:crop[1],crop[2]:crop[3]]
    else:
        imgCropped = img
    cv2.imwrite(f'{outputpath}/{figname}.png',imgCropped)
    return f'{outputpath}/{figname}.png'


def saveFigsPdf(pages, page, outputpath, figname, crop):
    with BytesIO() as image_byte_array:
        pages[page].save(image_byte_array, format='PNG')
        image_to_extract = image_byte_array.getvalue()
        img = cv2.imdecode(np.frombuffer(image_to_extract, np.uint8),3)
        if crop:
            imgCropped = img[crop[0]:crop[1],crop[2]:crop[3]]
        else:
            imgCropped = img
    cv2.imwrite(f'{outputpath}/{figname}.png',imgCropped)
    return f'{outputpath}/{figname}.png'

def saveFigsGeo(image, outputpath, figname, crop):
    if crop:
        imgCropped = image[crop[0]:crop[1],crop[2]:crop[3]]
    else:
        imgCropped = image
    image.save(f'{outputpath}/{figname}.png')
    return f'{outputpath}/{figname}.png'


def get_bold_list(para):
    bold_list = []
    for run in para.runs:
        bold_list.append(run.bold)
    return bold_list


def extractItemize_Imager(text):
    texs = """\\begin{itemize} \n """
    for i in text.split('•'):
        if len(i) > 1:
            texs += '\\item ' + ' '.join(i.split('\n')) +'\n'
            if i.startswith('TEC'):
                texs += '\\subsectoin{%s}' %(i) +'\n'
    
    texs += """\\end{itemize} \n """

    return texs


def extractTables(docum):
    tables = list()
    for table in docum.tables:

        data = []

        keys = None
        for i, row in enumerate(table.rows):
            text = (cell.text for cell in row.cells)

            if i == 0:
                keys = tuple(text)
                continue
            row_data = dict(zip(keys, text))
            data.append(row_data)
            # print (data)

        df = pd.DataFrame(data)
        tableStr = "\\begin{table}[H]\n\\centering\n \\begin{tabular}{|" + " | ".join(["c"] * len(df.columns)) + "|}\n"
        for i, row in df.iterrows():
            tableStr += "\\hline  \n"
            tableStr += " & ".join([str(x) for x in row.values]) + " \\\\\n"
        tableStr += "\\hline  \n\\end{tabular} \n \\caption{%s} \n \\end{table}"

        tables.append(tableStr)

    return tables


def separatePathsAreas(files):
    regexScint = r"(?i)\(Cintila.*S4\)"
    regexULF = r"(?i)ULF.*Geomagnetismo"
    regexSun0 = r"(?i)briefing\smeeting.*sun"
    regexSun1 = r"(?i)Resumo.*eventos.*solares"
    regexImager = r"(?i)LUME"
    regexEmic = r"(?i)ULF.*EMIC"
    regexROTI = r"(?i)ionosphere.*roti"
    regexRadBeltEn = r"(?i)radiation.*belt"
    regexRadBeltPt = r"(?i)cintur.*radia"
    regexIonosf = r"(?i)sum.*ionosfera"
    regexMeioInterpl = r"(?i)meio.*interpl"
    regexGeomag = r"(?i)(geomag)|(Embrace)"
    dictPaths = {}
    for ff in files:
        if ff.endswith("pdf"):
            with fitz.open(ff) as doc:
                text = ''
                for page in range(len(doc)):
                    text += doc[page].get_text()
            if re.search(regexULF, text, re.MULTILINE):
                dictPaths['ULF'] = {'path':ff}
            if re.search(regexScint, text, re.MULTILINE):
                dictPaths['Scintilation'] = {'path':ff}
            if re.search(regexSun0, text, re.MULTILINE):
                dictPaths['01Sun'] = {'path':ff}
            if re.search(regexSun1, text, re.MULTILINE):
                dictPaths['02Sun'] = {'path':ff}
            if re.search(regexImager, text, re.MULTILINE):
                dictPaths['Imager'] = {'path':ff}
            if re.search(regexEmic, text, re.MULTILINE):
                dictPaths['EMIC'] = {'path':ff}
        if ff.endswith("docx"):
            document = Document(ff)
            text= ''
            for para in document.paragraphs:
                text += para.text
            if re.search(regexROTI, text, re.MULTILINE):
                dictPaths['ROTI'] = {'path':ff}
            if re.search(regexRadBeltEn, text, re.MULTILINE):
                dictPaths['RadBelt'] = {'pathEn':ff}
            if re.search(regexRadBeltPt, text, re.MULTILINE):
                dictPaths['RadBelt'] = {'pathPt':ff}
            if re.search(regexIonosf, text, re.MULTILINE):
                dictPaths['Ionosfera'] = ff
        if ff.endswith("pptx"):
            prs = Presentation(ff)
            text= ''
            for ss in prs.slides:
                for shapes in ss.shapes:
                    if shapes.has_text_frame:
                        text += shapes.text
            
            if re.search(regexMeioInterpl, text, re.MULTILINE):
                dictPaths['03MeioInterp'] = ff
        if ff.endswith('xlsx'):
            wb = load_workbook(ff)
            sheetpt = wb['Plan4pt']
            text = ''
            for r in range(sheetpt.max_row):
                cell_obj = sheetpt.cell(row = r+1, column = 1)
                text += cell_obj.value
            if re.search(regexGeomag, text, re.MULTILINE):
                dictPaths['Geomag'] = ff
    
    return dictPaths